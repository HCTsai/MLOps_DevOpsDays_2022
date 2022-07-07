'''
Created on 2022年2月22日

@author: Hsiao-Chien Tsai
'''
from pptx import Presentation
import docx
import fitz # PyMuPdf
def docx_to_text(filename):
    doc = docx.Document(filename)
    fullText = []
    for para in doc.paragraphs:
        fullText.append(para.text)
    return "。".join(fullText)

def pptx_to_text(file_path):
    prs = Presentation(file_path)

    # text_runs will be populated with a list of strings,
    # one for each text run in presentation
    text_runs = []
    others = []
    
    for slide in prs.slides:
        for shape in slide.shapes:
            #print (shape.position)
            if not shape.has_text_frame:
                #print (shape)
                continue
            for paragraph in shape.text_frame.paragraphs:
                
                for run in paragraph.runs:
                    text_runs.append(run.text)
                #print (run.text)
            #text_runs.append("<shape>")
        #text_runs.append("<slide>")       
    res = "。".join(text_runs)
    return res
def pdf_to_text(file_path):
    doc = fitz.Document(file_path) 
    text = ""
    for page in doc:
        text += page.get_text() 
    return text
