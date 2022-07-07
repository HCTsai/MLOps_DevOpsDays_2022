'''
Created on 2022年2月14日

@author: Hsiao-Chien Tsai
'''

from pptx import Presentation
from pptx.util import Inches
from pptx.util import Cm, Pt
#from pptx.enum.text import PP_ALIGN
import matplotlib.pyplot as plt 
import time



def set_sowt_text(cell, strategy_type, strategy_dict, label_desc):
    p = cell.text_frame.add_paragraph()
    p.text = label_desc[strategy_type] + ":"
    strategies = strategy_dict[strategy_type]
    idx = 1 
    for k,v in strategies.items() :
        p = cell.text_frame.add_paragraph()
        p.text = "{}.{}".format(str(idx), k)
        p.font.size=Pt(14)
        idx += 1 
def set_pest_text(cell, strategy_type, strategy_dict, label_desc):
    strategies = strategy_dict[strategy_type]
    idx = 1
    for k,v in strategies.items() :
        p = cell.text_frame.add_paragraph()
        p.text = "{}.{}".format(str(idx), k)
        p.font.size=Pt(14)
        idx += 1 
        
def create_swot_slide(prs, strategy_dict, label_desc, pptx_title):
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes
    shapes.title.text = pptx_title
    shapes.title.text_frame.paragraphs[0].font.size=Pt(28)
    
    rows = 3
    cols = 2
    left = Inches(0.3)
    top = Inches(1.2)
    width = Inches(9.0)
    height = Inches(0.8)
    table = shapes.add_table(rows, cols, left, top, width, height).table

    # set column widths
    table.columns[0].width = Inches(4.5)
    table.columns[1].width = Inches(4.5)
    # write column headings
    table.cell(0, 0).text = "Helpful"
    table.cell(0, 1).text = "Harmful"
   
    # write body cells
    
    set_sowt_text(table.cell(1, 0), "S", strategy_dict, label_desc)
    
    '''
    paragraph = cell.text_frame.paragraphs[0]
    run = paragraph.add_run() 
    run.text = "123"
    '''
    set_sowt_text(table.cell(1, 1), "W", strategy_dict, label_desc)
    set_sowt_text(table.cell(2, 0), "O", strategy_dict, label_desc)
    set_sowt_text(table.cell(2, 1), "T", strategy_dict, label_desc)
    return prs
def create_pest_slide(prs, strategy_dict, label_desc, pptx_title):
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes
    shapes.title.text = pptx_title
    shapes.title.text_frame.paragraphs[0].font.size=Pt(28)
    
    rows = 2
    cols = 4
    left = Inches(0.3)
    top = Inches(1.2)
    width = Inches(9.0)
    height = Inches(6)
    table = shapes.add_table(rows, cols, left, top, width, height).table

    # set column widths
    table.rows[0].height = Inches(0.5)
    #table.columns[1].width = Inches(4.5)
    # write column headings
    #table.cell(0, 0).text = ''
    #table.cell(0, 1).text = ''
    # write body cells
    idx = 0 
    for k,v in label_desc.items() :
        #title
        table.cell(0, idx).text = v
        table.cell(0, idx).text_frame.paragraphs[0].font.size=Pt(16)
        #strategy
        set_pest_text(table.cell(1, idx), k, strategy_dict, label_desc)
        idx += 1
    return prs
def save_strategy_to_pptx(strategy_dict, label_desc, pptx_title, out_file_name, graph_type) :
    prs = Presentation()
    if (graph_type == "SWOT") :
        prs = create_swot_slide(prs, strategy_dict, label_desc, pptx_title)
    if (graph_type == "PEST" or graph_type == "DISC") :
        prs = create_pest_slide(prs, strategy_dict, label_desc, pptx_title)
        
    
    prs.save(out_file_name)

